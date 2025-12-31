# -*- coding: utf-8 -*-
import pandas as pd
import tkinter as tk
from tkinter import filedialog, messagebox
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------
# Chart Styling Helper
# ---------------------------
def style_chart(chart, title_text):
    chart.set_plotarea({'fill': {'color': '#f4f6f8'}})
    chart.set_chartarea({'fill': {'color': '#e9eff5'}})
    chart.set_title({
        'name': title_text,
        'name_font': {'bold': True, 'color': '#1f4e79', 'size': 11}
    })

# ---------------------------
# PPT HELPERS
# ---------------------------
def set_slide_background(slide, rgb_color=(230, 240, 250)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)


def _add_header(slide, headline_text):
    # Blue header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10.00),
        Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(79, 121, 191)
    header.line.fill.background()

    # Headline text (left)
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.2),
        Inches(9),
        Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = headline_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Jio logo (right – vector, no image dependency)
    logo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.2),
        Inches(0.10),
        Inches(0.7),
        Inches(0.7)
    )
    logo.fill.solid()
    logo.fill.fore_color.rgb = RGBColor(0, 0, 128)
    logo.line.fill.background()

    logo_tf = logo.text_frame
    lp = logo_tf.paragraphs[0]
    lp.text = "Jio"
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(255, 255, 255)
    lp.alignment = 1  # center

def _add_bar_chart(slide, title, categories, series):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.3), Inches(9.0), Inches(4.5),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(9)

    # =========================
    # ADD VALUES ON BARS (NEW)
    # =========================
    for plot in chart.plots:
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_value = True
        data_labels.font.size = Pt(9)


def _add_pie_chart(slide, title, labels, values):
    data_pairs = sorted(zip(labels, values), key=lambda x: x[1], reverse=True)
    if len(data_pairs) > 10:
        top10 = data_pairs[:10]
        others = data_pairs[10:]
        labels = [l for l,_ in top10] + ["Others"]
        values = [v for _,v in top10] + [sum(v for _,v in others)]
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(title, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1.5), Inches(1.3), Inches(6.0), Inches(4.8),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(9)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True

# ---------------------------
# PPT GENERATION
# ---------------------------
def generate_ppt(ppt_path, category_summary, category_business, eta_business, daily_eta, all_issues, df):
    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    _add_header(slide, "Total Alerts per Category")
    _add_bar_chart(slide, "", category_summary['category'].astype(str).tolist(),
                   [('Total Alerts', category_summary['total_issues'].astype(int).tolist())])

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    _add_header(slide, "Business vs Non-Business by Category")
    _add_bar_chart(slide, "",
                   category_business['category'].astype(str).tolist(),
                   [('Business', category_business['Business'].tolist()),
                    ('Non-Business', category_business['Non-Business'].tolist())])

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    _add_header(slide, "ETA Breach – Business vs Non-Business")
    yes_row = eta_business[eta_business['ETA_Breach'] == 'Yes'].index[0]
    _add_bar_chart(slide, "",
                   ['Business', 'Non-Business'],
                   [('Count', [
                       int(eta_business.loc[yes_row, 'Business']),
                       int(eta_business.loc[yes_row, 'Non-Business'])
                   ])])

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    _add_header(slide, "Daily ETA Breach Trend")
    _add_bar_chart(slide, "",
                   daily_eta['createdOn'].astype(str).tolist(),
                   [('Breached Count', daily_eta['Breached_Count'].tolist())])

    # Application Slides
    for app, app_df in all_issues.groupby('applicationName'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        _add_header(slide, f"Top 10 Issues – {app}")
        top10 = app_df.head(10)
        if not top10.empty:
            _add_pie_chart(slide, "",
                           top10['example_value'].astype(str).str[:30].tolist(),
                           top10['repeat_count'].tolist())

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        _add_header(slide, f"Business vs Non-Business – {app}")
        biz = df[df['applicationName'] == app].groupby('Business').size().reindex(
            ['Business', 'Non-Business'], fill_value=0)
        _add_bar_chart(slide, "", ['Business', 'Non-Business'], [('Count', biz.tolist())])

    prs.save(ppt_path)

# ---------------------------
# Main Logic
# ---------------------------
def process_excel():
    file_path = filedialog.askopenfilename(title="Select Excel file", filetypes=[("Excel files","*.xlsx *.xls")])
    if not file_path:
        return
    try:
        df = pd.read_excel(file_path)
        logic_col = 'kpiName' if use_kpi_var.get() else 'title'
        clean_col = f"{logic_col}_clean"
        required_cols = ['ipAddress', logic_col, 'category','applicationName','ackMetStatus','createdOn','environment','latestUpdCategory']
        missing = [c for c in required_cols if c not in df.columns]
        if missing:
            messagebox.showerror("Error", f"Missing columns: {missing}")
            return

        # Normalize Title / KPI
        df[clean_col] = (
            df[logic_col].astype(str).str.lower()
            .apply(lambda x: re.sub(r'\d+', '<num>', x))
            .str.replace(r'[^a-z0-9 <>]', ' ', regex=True)
            .str.replace(r'\s+', ' ', regex=True)
            .str.strip()
        )

        # ETA Logic
        df['ackMetStatus'] = df['ackMetStatus'].astype(str).str.upper().str.strip()
        df['ETA_Breach'] = df['ackMetStatus'].map({'Y':'No','N':'Yes','X':'Yes'}).fillna('Yes')
        df['createdOn'] = pd.to_datetime(df['createdOn'])
        df['hour'] = df['createdOn'].dt.hour
        df['environment'] = df['environment'].astype(str).str.upper()
        df['latestUpdCategory'] = df['latestUpdCategory'].astype(str).str.upper()

        business_hour = (df['hour'] >=7) & (df['hour'] <23)
        exclude_env = df['environment'].str.contains('DR|REPLICA', regex=True)
        exclude_cr = df['latestUpdCategory'].str.contains('SUPPRESSED BY CR', regex=True)

        df['Business'] = 'Non-Business'
        df.loc[business_hour & ~exclude_env & ~exclude_cr, 'Business'] = 'Business'
        df.drop(columns=['hour'], inplace=True)

        # Issue aggregation
        issue_counts = (
            df.groupby(['ipAddress','applicationName','category',clean_col])
            .agg(repeat_count=(clean_col,'size'), example_value=(logic_col,'first'))
            .reset_index()
        )
        all_issues = issue_counts[['ipAddress','applicationName','category','example_value','repeat_count',clean_col]].sort_values('repeat_count',ascending=False)

        # Category summaries
        category_summary = issue_counts.groupby('category').agg(total_issues=('repeat_count','sum'),unique_alerts=(clean_col,'nunique')).reset_index().sort_values('total_issues',ascending=False)
        category_business = df.groupby(['category','Business']).size().unstack(fill_value=0).reindex(columns=['Business','Non-Business'],fill_value=0).reset_index()
        eta_summary = df.groupby('ETA_Breach').size().reindex(['Yes','No'],fill_value=0).reset_index(name='count')
        eta_business = df.groupby(['ETA_Breach','Business']).size().unstack(fill_value=0).reindex(columns=['Business','Non-Business'],fill_value=0).reset_index()
        daily_eta = df[df['ETA_Breach']=='Yes'].groupby(df['createdOn'].dt.date).size().reset_index(name='Breached_Count')
        daily_eta['createdOn'] = daily_eta['createdOn'].astype(str)

        output_path = filedialog.asksaveasfilename(defaultextension=".xlsx",filetypes=[("Excel files","*.xlsx *.xls")])
        if not output_path:
            return

        with pd.ExcelWriter(output_path, engine='xlsxwriter') as writer:
            workbook = writer.book
            # ====== ALL Excel Sheets & Charts Logic (same as mother code) ======
            # All_Alert
            all_issues.to_excel(writer,index=False,sheet_name='All_Alert')
            # All_Alert_Business
            all_issues_business = df.groupby(['ipAddress','applicationName','category',clean_col,'Business'])\
                                    .agg(repeat_count=(clean_col,'size'),example_value=(logic_col,'first'))\
                                    .reset_index().sort_values('repeat_count',ascending=False)
            all_issues_business.to_excel(writer,index=False,sheet_name='All_Alert_Business')
            all_alert_ws = writer.sheets['All_Alert_Business']
            biz_count = all_issues_business.groupby('Business')['repeat_count'].sum().reindex(['Business','Non-Business'],fill_value=0).reset_index()
            start_biz_chart = len(all_issues_business)+3
            biz_count.to_excel(writer,sheet_name='All_Alert_Business',startrow=start_biz_chart,index=False)
            all_alert_chart = workbook.add_chart({'type':'column'})
            all_alert_chart.add_series({
                'name':'Business vs Non-Business',
                'categories':['All_Alert_Business',start_biz_chart+1,0,start_biz_chart+2,0],
                'values':['All_Alert_Business',start_biz_chart+1,1,start_biz_chart+2,1],
                'points':[{'fill':{'color':'red'}},{'fill':{'color':'green'}}],
                'data_labels':{'value':True}
            })
            style_chart(all_alert_chart,'All Alerts – Business vs Non-Business')
            all_alert_ws.insert_chart('H2',all_alert_chart)

            # Category_Summary
            category_summary.to_excel(writer,index=False,sheet_name='Category_Summary')
            cat_ws = writer.sheets['Category_Summary']
            cat_chart = workbook.add_chart({'type':'column'})
            cat_chart.add_series({
                'name':'Total Alerts',
                'categories':['Category_Summary',1,0,len(category_summary),0],
                'values':['Category_Summary',1,1,len(category_summary),1],
                'data_labels':{'value':True}
            })
            style_chart(cat_chart,'Total Alerts per Category')
            cat_ws.insert_chart('E2',cat_chart)
            start = len(category_summary)+4
            category_business.to_excel(writer,sheet_name='Category_Summary',startrow=start,index=False)
            cat_biz = workbook.add_chart({'type':'column'})
            cat_biz.add_series({
                'name':'Business',
                'categories':['Category_Summary',start+1,0,start+len(category_business),0],
                'values':['Category_Summary',start+1,1,start+len(category_business),1],
                'fill':{'color':'red'},
                'data_labels':{'value':True}
            })
            cat_biz.add_series({
                'name':'Non-Business',
                'categories':['Category_Summary',start+1,0,start+len(category_business),0],
                'values':['Category_Summary',start+1,2,start+len(category_business),2],
                'fill':{'color':'green'},
                'data_labels':{'value':True}
            })
            style_chart(cat_biz,'Business vs Non-Business by Category')
            cat_ws.insert_chart('E20',cat_biz)

            # =========================
            # NEW SHEET: Application vs Category Matrix
            # =========================
            app_category_matrix = (
            issue_counts
            .pivot_table(
            index='applicationName',
            columns='category',
            values='repeat_count',
            aggfunc='sum',
            fill_value=0)
            )

            # Add Grand Total column
            app_category_matrix['Grand Total'] = app_category_matrix.sum(axis=1)

            # Add Grand Total row
            grand_total_row = pd.DataFrame(
            app_category_matrix.sum(axis=0)
            ).T
            grand_total_row.index = ['Grand Total']

            app_category_matrix = pd.concat([app_category_matrix, grand_total_row])

                # Reset index for Excel
            app_category_matrix.reset_index(inplace=True)
            app_category_matrix.rename(columns={'applicationName': 'Application Name'}, inplace=True)
            app_category_matrix.rename(columns={'index': 'Application Name'}, inplace=True)
                
            app_category_matrix.to_excel(
                writer,
                sheet_name='Application_Category_Matrix',
                index=False
                 )
               

            # ETA_Ack_Breach
            eta_summary.to_excel(writer,index=False,sheet_name='ETA_Ack_Breach')
            eta_ws = writer.sheets['ETA_Ack_Breach']
            eta_chart = workbook.add_chart({'type':'column'})
            eta_chart.add_series({
                'name':'ETA Status',
                'categories':['ETA_Ack_Breach',1,0,2,0],
                'values':['ETA_Ack_Breach',1,1,2,1],
                'points':[{'fill':{'color':'red'}},{'fill':{'color':'green'}}],
                'data_labels':{'value':True}
            })
            style_chart(eta_chart,'ACK-ETA Breach Status')
            eta_ws.insert_chart('D2',eta_chart)
            start_eta = len(eta_summary)+4
            eta_business.to_excel(writer,sheet_name='ETA_Ack_Breach',startrow=start_eta,index=False)
            eta_yes_row = eta_business[eta_business['ETA_Breach']=='Yes'].index[0]
            excel_row = start_eta + 1 + eta_yes_row
            eta_biz = workbook.add_chart({'type':'column'})
            eta_biz.add_series({
                'name':'Business',
                'categories':['ETA_Ack_Breach',excel_row,0,excel_row,0],
                'values':['ETA_Ack_Breach',excel_row,1,excel_row,1],
                'fill':{'color':'red'},
                'data_labels':{'value':True}
            })
            eta_biz.add_series({
                'name':'Non-Business',
                'categories':['ETA_Ack_Breach',excel_row,0,excel_row,0],
                'values':['ETA_Ack_Breach',excel_row,2,excel_row,2],
                'fill':{'color':'green'},
                'data_labels':{'value':True}
            })
            style_chart(eta_biz,'ACK-ETA Breach – Business vs Non-Business')
            eta_ws.insert_chart('D20',eta_biz)
            trend_start = start_eta + len(eta_business)+6
            daily_eta.to_excel(writer,sheet_name='ETA_Ack_Breach',startrow=trend_start,index=False)
            trend_chart = workbook.add_chart({'type':'line'})
            trend_chart.add_series({
                'name':'Daily ACK-ETA Breach Trend',
                'categories':['ETA_Ack_Breach',trend_start+1,0,trend_start+len(daily_eta),0],
                'values':['ETA_Ack_Breach',trend_start+1,1,trend_start+len(daily_eta),1],
                'data_labels':{'value':True},
                'marker':{'type':'circle'}
            })
            trend_chart.set_x_axis({'text_axis':True,'label_position':'low'})
            style_chart(trend_chart,'Daily ACK-ETA Breach Trend')
            eta_ws.insert_chart('K2',trend_chart)

            # ETA_Alert
            df.to_excel(writer,index=False,sheet_name='ETA_Alert')
            
            # Application Sheets
            for app, app_df in all_issues.groupby('applicationName'):
                sheet = f"App_{app}"[:31]
                app_df.to_excel(writer,index=False,sheet_name=sheet)
                ws = writer.sheets[sheet]
                top10 = app_df.head(10)
                if not top10.empty:
                    top_chart = workbook.add_chart({'type':'bar'})
                    top_chart.add_series({
                        'name':'Repeat Count',
                        'categories':[sheet,1,3,len(top10),3],
                        'values':[sheet,1,4,len(top10),4],
                        'data_labels':{'value':True}
                    })
                    style_chart(top_chart,f'Top 10 Alerts – {app}')
                    ws.insert_chart('H2',top_chart)
                biz = df[df['applicationName']==app].groupby('Business').size().reindex(['Business','Non-Business'],fill_value=0).reset_index(name='count')
                start_app = len(app_df)+4
                biz.to_excel(writer,sheet_name=sheet,startrow=start_app,index=False)
                app_chart = workbook.add_chart({'type':'column'})
                app_chart.add_series({
                    'name':'Business vs Non-Business',
                    'categories':[sheet,start_app+1,0,start_app+2,0],
                    'values':[sheet,start_app+1,1,start_app+2,1],
                    'points':[{'fill':{'color':'red'}},{'fill':{'color':'green'}}],
                    'data_labels':{'value':True}
                })
                style_chart(app_chart,f'Business vs Non-Business – {app}')
                ws.insert_chart('H20',app_chart)
                
        # Generate PPT
        ppt_path = output_path.replace('.xlsx','.pptx')
        generate_ppt(ppt_path, category_summary, category_business, eta_business, daily_eta, all_issues, df)
        messagebox.showinfo("Success",f"Report generated successfully.\nExcel: {output_path}\nPPT: {ppt_path}")

    except Exception as e:
        messagebox.showerror("Error", str(e))

# ---------------------------
# GUI
# ---------------------------
root = tk.Tk()
root.title("Repeated NGO-Alert Report")
root.geometry("760x420")

use_kpi_var = tk.BooleanVar(value=False)

tk.Label(root,text="Alert Grouping Logic",font=("Arial",13,"bold")).pack(pady=10)
tk.Checkbutton(root,text="Use KPI Name instead of Title",variable=use_kpi_var,font=("Arial",11)).pack()
tk.Button(root,text="Select Excel and Process",font=("Arial",12),command=process_excel).pack(pady=25)

root.mainloop()